"""
Script de Generación de Diapositivas para Divulgación Científica.
CORREGIDO: Genera títulos manuales para evitar errores en layouts vacíos.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- CONFIGURACIÓN DE ESTILO (PALETA QUANTUM) ---
# Fondo: Azul Noche Profundo (Professional & Tech)
BG_COLOR = RGBColor(11, 16, 33)  # #0B1021
# Acento: Cian Eléctrico (Para títulos y énfasis)
ACCENT_COLOR = RGBColor(0, 255, 255) # Cyan
# Texto Principal: Blanco Hueso (Para lectura cómoda)
TEXT_COLOR = RGBColor(240, 240, 240)
# Texto Secundario: Gris Plata
SUBTEXT_COLOR = RGBColor(180, 180, 190)

def set_background(slide):
    """Aplica el modo oscuro al slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title(slide, text, font_size=44):
    """
    Añade un título manualmente (creando un TextBox) para evitar
    errores en diapositivas en blanco (Layout 6).
    """
    # Definir posición y tamaño del título
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1.5)
    
    # Crear el cuadro de texto explícitamente
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    # Configurar el texto y estilo
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = "Arial" # Tipografía segura y limpia
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.LEFT

def add_content_text(slide, text_list, top_margin=1.5):
    """Añade lista de puntos con estilo limpio."""
    left = Inches(0.5)
    top = Inches(top_margin)
    width = Inches(9)
    height = Inches(5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(text_list):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(24) if i == 0 else Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = "Arial"
        p.space_after = Pt(14)
        
        # Simular viñetas modernas con caracteres ASCII si no es el primer párrafo
        if i > 0:
            p.text = "  • " + line

def generate_presentation():
    print("🎨 Generando presentación 'Dark Quantum'...")
    prs = Presentation()
    
    # 1. SLIDE DE PORTADA (IMPACTO)
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Layout vacío
    set_background(slide)
    
    # Título Central
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Detección de Fraude con\nInteligencia Artificial Cuántica"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.add_paragraph()
    p_sub.text = "Un enfoque híbrido usando Qiskit, CUDA y Hamiltonianos de Ising"
    p_sub.font.size = Pt(24)
    p_sub.font.color.rgb = TEXT_COLOR
    p_sub.alignment = PP_ALIGN.CENTER

    # 2. SLIDE DEL PROBLEMA (CONTEXTO)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "El Desafío: Fraude Sofisticado")
    
    content = [
        "¿Por qué fallan los modelos clásicos?",
        "Los métodos tradicionales (RBF) asumen fronteras suaves.",
        "El fraude moderno se camufla: son anomalías que imitan transacciones reales.",
        "Necesitamos mayor expresividad matemática para separar los datos.",
        "Solución Propuesta: Mapear datos a un Espacio de Hilbert (2^N dimensiones)."
    ]
    add_content_text(slide, content)

    # 3. SLIDE DE TEORÍA (EL CEREBRO CUÁNTICO)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "La Física: Hamiltoniano de Ising")
    
    # Explicación
    content_text = [
        "Usamos un circuito variacional que evoluciona los datos bajo interacción ZZ.",
        "Esto induce 'Entrelazamiento Cuántico' controlado por los datos."
    ]
    add_content_text(slide, content_text, top_margin=1.2)
    
    # Simulación de ecuación
    eq_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1.5))
    tf_eq = eq_box.text_frame
    p_eq = tf_eq.add_paragraph()
    p_eq.text = "H(x) = Σ Zᵢ + Σ ZᵢZⱼ (Interacción Spin-Spin)"
    p_eq.font.size = Pt(32)
    p_eq.font.name = "Courier New"
    p_eq.font.bold = True
    p_eq.font.color.rgb = RGBColor(255, 0, 255) # Magenta
    p_eq.alignment = PP_ALIGN.CENTER
    
    # Nota al pie
    note_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    tf_note = note_box.text_frame
    p_note = tf_note.add_paragraph()
    p_note.text = "Las correlaciones no lineales permiten detectar patrones invisibles para la estadística clásica."
    p_note.font.size = Pt(18)
    p_note.font.color.rgb = SUBTEXT_COLOR
    p_note.alignment = PP_ALIGN.CENTER

    # 4. SLIDE DE EVIDENCIA (TU FIGURA 1)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Topología de Decisión: La Prueba Visual")
    
    img_path = os.path.join("references", "Figure_1.png")
    # Intentar cargar png o pdf (python-pptx prefiere png/jpg)
    if not os.path.exists(img_path):
         # fallback por si se guardó con otro nombre
         img_path = os.path.join("references", "figure_1_decision_boundary.png")

    if os.path.exists(img_path):
        # Insertar imagen centrada
        pic = slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), height=Inches(5))
        
        # Añadir "Callout"
        comm_box = slide.shapes.add_textbox(Inches(6.5), Inches(2), Inches(3), Inches(4))
        tf_comm = comm_box.text_frame
        p_comm = tf_comm.add_paragraph()
        p_comm.text = "Observa la diferencia:"
        p_comm.font.bold = True
        p_comm.font.size = Pt(20)
        p_comm.font.color.rgb = ACCENT_COLOR
        
        p_list = tf_comm.add_paragraph()
        p_list.text = "\nPanel Izq (Clásico):\nFronteras redondas y suaves.\n\nPanel Der (Cuántico):\n'Islas' de decisión fragmentadas.\n\nEl modelo cuántico aísla el fraude con precisión quirúrgica."
        p_list.font.size = Pt(16)
        p_list.font.color.rgb = TEXT_COLOR
    else:
        print(f"⚠️ ALERTA: No se encontró la imagen en {img_path}. Genera la gráfica primero.")

    # 5. SLIDE DE RESULTADOS (LA TABLA)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Resultados del Benchmark")
    
    # Crear Tabla
    rows, cols = 3, 3
    left, top, width, height = Inches(1.5), Inches(2.5), Inches(7), Inches(2)
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table
    
    # Headers
    headers = ["Modelo", "Accuracy", "Observación"]
    data = [
        ["SVM Clásico (RBF)", "97.6%", "Rápido y Generalista"],
        ["SVM Cuántico (Ising)", "95.2%", "Alta Expresividad"]
    ]
    
    # Estilo Headers
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(50, 50, 80)
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.color.rgb = ACCENT_COLOR
        para.font.size = Pt(20)

    # Estilo Datos
    for row_idx, row_data in enumerate(data):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(20, 25, 45)
            para = cell.text_frame.paragraphs[0]
            para.font.color.rgb = TEXT_COLOR
            para.font.size = Pt(18)

    # Conclusión abajo
    conc_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    tf_c = conc_box.text_frame
    p_c = tf_c.add_paragraph()
    p_c.text = "Conclusión: El modelo cuántico es competitivo y ofrece una topología superior para detectar fraudes complejos."
    p_c.font.size = Pt(18)
    p_c.font.color.rgb = SUBTEXT_COLOR
    p_c.alignment = PP_ALIGN.CENTER

    # 6. SLIDE DE CIERRE (CTA)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    final_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    tf_f = final_box.text_frame
    
    p_f = tf_f.add_paragraph()
    p_f.text = "Gracias"
    p_f.font.size = Pt(60)
    p_f.font.bold = True
    p_f.font.color.rgb = ACCENT_COLOR
    p_f.alignment = PP_ALIGN.CENTER
    
    p_link = tf_f.add_paragraph()
    p_link.text = "\nCódigo y Paper disponible en GitHub\n#Qiskit #DataScience #QuantumComputing"
    p_link.font.size = Pt(20)
    p_link.font.color.rgb = TEXT_COLOR
    p_link.alignment = PP_ALIGN.CENTER

    # Guardar
    if not os.path.exists("references"):
        os.makedirs("references")
    output_path = "references/Presentacion_Divulgacion_Quantum.pptx"
    prs.save(output_path)
    print(f"✅ Presentación generada exitosamente en: {output_path}")

if __name__ == "__main__":
    generate_presentation()